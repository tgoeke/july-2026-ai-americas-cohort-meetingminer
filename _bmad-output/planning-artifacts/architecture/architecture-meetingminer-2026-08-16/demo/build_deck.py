# Builds meetingminer-cohort-demo.pptx from the diagrams/ PNGs.
# Regenerate: uv run --with python-pptx python build_deck.py
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
OUT = HERE / "meetingminer-cohort-demo.pptx"
DIA = HERE / "diagrams"

ACCENT = RGBColor(0x1E, 0x5A, 0x8A)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x4A, 0x55, 0x68)
LIGHT = RGBColor(0xE8, 0xF0, 0xF8)
FONT = "Helvetica Neue"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, left, top, width, height):
    tb = s.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def para(tf, text, size, *, bold=False, color=DARK, first=False, bullet=False,
         space_after=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    f = r.font
    f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, color
    return p


def header(s, title, timebox):
    bar = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tb = box(s, Inches(0.5), Inches(0.25), Inches(10.8), Inches(0.7))
    para(tb.text_frame, title, 30, bold=True, color=ACCENT, first=True)
    tt = box(s, Inches(11.4), Inches(0.33), Inches(1.6), Inches(0.5))
    para(tt.text_frame, timebox, 14, color=GRAY, first=True, align=PP_ALIGN.RIGHT)


def picture(s, path, left, top, max_w, max_h):
    from PIL import Image
    w, h = Image.open(path).size
    scale = min(max_w / w, max_h / h)
    pw, ph = int(w * scale), int(h * scale)
    return s.shapes.add_picture(str(path), left + Emu(int((max_w - pw) / 2)),
                                top + Emu(int((max_h - ph) / 2)), pw, ph)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ── 1 · Title ────────────────────────────────────────────────────────────────
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.5), SLIDE_W, Inches(2.6))
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()
tb = box(s, Inches(1), Inches(2.75), Inches(11.3), Inches(2.1))
para(tb.text_frame, "MeetingMiner", 54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), first=True)
para(tb.text_frame, "Evidence-first meeting intelligence — search your meetings, replay the exact moment, "
     "trust every answer.", 22, color=LIGHT)
tb = box(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.9))
para(tb.text_frame, "Cameron Blake  ·  AI Cohort Demo  ·  August 22, 2026", 16, color=GRAY, first=True)
notes(s, "0:00–0:05 — MeetingMiner turns Teams meetings into searchable, replayable evidence. "
         "Three minutes: the architecture, the decisions, then it live.")

# ── 2 · What it does ─────────────────────────────────────────────────────────
s = slide()
header(s, "What it does", "0:05–0:15")
cols = [
    ("INGEST", "Teams recordings + transcripts land as write-once evidence drops. "
               "A checkpointed pipeline derives frames, screens, transcripts, moments."),
    ("SEARCH", "Hybrid keyword + semantic search over every meeting moment — "
               "snippets highlighted, filterable by meeting and corpus."),
    ("REPLAY", "Every result is a moment with millisecond timing — click it and the "
               "video plays from that exact point. Answers cite evidence you can watch."),
]
for i, (t, body) in enumerate(cols):
    left = Inches(0.5 + i * 4.35)
    card = s.shapes.add_shape(1, left, Inches(1.7), Inches(4.0), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = ACCENT
    tb = box(s, left + Inches(0.25), Inches(1.95), Inches(3.5), Inches(3.1))
    para(tb.text_frame, t, 22, bold=True, color=ACCENT, first=True, space_after=10)
    para(tb.text_frame, body, 15, color=DARK)
tb = box(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.9))
para(tb.text_frame, "Principle: no citation, no answer — every claim traces to a replayable moment.",
     18, bold=True, color=ACCENT, first=True, align=PP_ALIGN.CENTER)
notes(s, "0:05–0:15 — Meetings are where decisions happen and where they vanish. MeetingMiner ingests "
         "them once, then everything is search-and-replay. The rule that shapes the whole design: "
         "no citation, no answer.")

# ── 3 · Finding the screenshots ─────────────────────────────────────────────
s = slide()
header(s, "Finding the screenshots", "0:15–0:30")
picture(s, DIA / "screens.png", Inches(0.4), Inches(1.3), Inches(12.5), Inches(1.6))
tb = box(s, Inches(0.7), Inches(3.2), Inches(11.9), Inches(3.9))
para(tb.text_frame, "A screen is identified by what it says, not how it looks — OCR-text "
     "similarity survives resolution, compression, and webcam overlays where pixel matching breaks",
     16, bullet=True, first=True, space_after=10)
para(tb.text_frame, "Identity is cross-meeting: the same deck shown next week resolves to the same "
     "screen — the lineage of a slide across the whole corpus", 16, bullet=True, space_after=10)
para(tb.text_frame, "Every screenshot records its frame span and capture cue; moments cite it — "
     "the visual half of every replayable citation", 16, bullet=True, space_after=10)
para(tb.text_frame, "Tuned on the measured corpus for slides and app views; dense spreadsheet-style "
     "screens are accepted misses", 16, bullet=True)
notes(s, "0:15–0:30 — Before any AI answers anything, we find what was on screen. A frame every two "
         "seconds, crop the webcam column, OCR each frame — then screen identity comes from the text, "
         "not the pixels. That's what lets the same slide match across meetings, and it gives every "
         "citation its screenshot.")

# ── 4 · Architecture ─────────────────────────────────────────────────────────
s = slide()
header(s, "Architecture", "0:30–0:45")
picture(s, DIA / "architecture.png", Inches(0.3), Inches(1.15), Inches(8.9), Inches(6.1))
tb = box(s, Inches(9.4), Inches(1.5), Inches(3.6), Inches(5.6))
para(tb.text_frame, "The paradigm", 18, bold=True, color=ACCENT, first=True, space_after=8)
para(tb.text_frame, "Deterministic evidence pipeline — model output never writes evidence", 14, bullet=True)
para(tb.text_frame, "Ports & adapters — every model call goes through a config-bound port", 14, bullet=True)
para(tb.text_frame, "CQRS-lite — Postgres is the record; Neo4j + Meilisearch are disposable "
     "projections with exactly one writer", 14, bullet=True)
para(tb.text_frame, "Infra in Docker, code on the Mac host — pipeline needs Apple Vision + MLX/Metal",
     14, bullet=True)
notes(s, "0:30–0:45 — One paradigm sentence carries it: deterministic evidence pipeline, ports-and-adapters "
         "at every model boundary, CQRS-lite storage. Left to right: the puller acquires from Teams, the "
         "worker derives evidence into Postgres, and projects it into a graph store and a search store. "
         "Only deterministic code writes evidence; models are behind swappable ports.")

# ── 5 · Technology stack — and why ──────────────────────────────────────────
s = slide()
header(s, "Technology stack — and why", "0:45–1:00")
rows = [
    ("Layer", "Choice", "Why"),
    ("API + worker", "Python 3.12 · FastAPI 0.141", "Typed, SSE streaming, OpenAPI → generated TS client"),
    ("Web", "React 19 · Vite 8 · Tailwind 4 + shadcn", "Fast SPA; client generated from the live API schema"),
    ("Record", "Postgres 18", "Single database of record; native UUIDv7 minting"),
    ("Graph", "Neo4j Community 2026.07", "Deterministic Cypher traversals over the evidence graph"),
    ("Search", "Meilisearch 1.53", "Hybrid keyword + vector in one store; 1024-dim user-provided embeddings"),
    ("Inference", "MLX Whisper · Apple Vision OCR · claude-sonnet-5 (Ollama fallback)",
     "Local-first where quality allows; every engine swappable in config.yaml — never a code change"),
    ("Runtime", "Docker for the 3 stores only; code as macOS host processes",
     "Pipeline needs Metal/MLX + Apple frameworks containers can't reach"),
]
top = Inches(1.35)
widths = [Inches(1.9), Inches(4.6), Inches(5.9)]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.5), top, sum(widths, Emu(0)), Inches(5.4)).table
for i, w in enumerate(widths):
    tbl.columns[i].width = w
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.name, p.font.size = FONT, Pt(13 if r else 14)
        p.font.bold = r == 0 or c == 0
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else DARK
notes(s, "0:45–1:00 — Boring, current technology on purpose. The one unusual call: Docker runs only the "
         "stateful stores — the pipeline runs on the Mac host because it needs Apple Vision and MLX. "
         "Every model — STT, OCR, LLM, embedder — sits behind a port bound in one config file; swapping "
         "a model is a config edit, never code.")

# ── 6 · Storage & RAG decisions ─────────────────────────────────────────────
s = slide()
header(s, "The critical storage + RAG decisions", "1:00–1:15")
tb = box(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.6))
items = [
    ("One database of record.", "Every entity is a Postgres row first; its UUIDv7 is its identity in "
     "every store, payload, and citation. No second owner exists."),
    ("Projections have exactly one writer.", "All Neo4j + Meilisearch writes go through one module; "
     "both stores rebuild from Postgres + config.yaml alone (make rebuild). Corruption is answered by "
     "rebuild, never by hand-editing an index."),
    ("GraphRAG is deterministic traversal.", "Retrieval is hand-written, parameterized Cypher templates "
     "+ hybrid search. The LLM only classifies the question and writes the answer — it never owns "
     "retrieval, so retrieval is testable."),
    ("No citation, no answer — enforced in code.", "A validator resolves every [[moment:uuid]] marker "
     "against Postgres; an uncited claim means the answer never leaves the API. Not a prompt instruction."),
    ("Drafts never reach retrieval.", "Extracted artifacts are human-approved before publishing; the "
     "publish gate lives inside the projection module."),
]
first = True
for head, body in items:
    p = para(tb.text_frame, head, 17, bold=True, color=ACCENT, first=first, space_after=2)
    first = False
    para(tb.text_frame, body, 14, color=DARK, space_after=10)
notes(s, "1:00–1:15 — Five decisions carry the design. Postgres is the only record. The two retrieval "
         "stores are disposable projections with a single writer — we can delete and rebuild them from "
         "the record. GraphRAG here means deterministic Cypher, not a framework owning a graph. And the "
         "citation rule is a code gate, not a prompt.")

# ── 7 · Encodings ────────────────────────────────────────────────────────────
s = slide()
header(s, "How evidence is encoded", "1:15–1:25")
tb = box(s, Inches(0.5), Inches(1.35), Inches(6.1), Inches(5.6))
para(tb.text_frame, "Embeddings", 18, bold=True, color=ACCENT, first=True, space_after=4)
para(tb.text_frame, "qwen3-embedding 0.6b (Ollama), fixed 1024-dim, computed by the projection module — "
     "store auto-embedders stay off", 14, bullet=True)
para(tb.text_frame, "Chunks: 1,400 chars with 1-turn overlap; hybrid ratio 0.3 — keyword-heavy on purpose",
     14, bullet=True)
para(tb.text_frame, "Model + dimension are projection state: a mismatch fails closed; changing the model "
     "forces a full rebuild + eval rerun", 14, bullet=True)
tb = box(s, Inches(6.9), Inches(1.35), Inches(6.0), Inches(5.6))
para(tb.text_frame, "Transcripts & video", 18, bold=True, color=ACCENT, first=True, space_after=4)
para(tb.text_frame, "Provided transcripts are immutable; the STT verification lane (MLX whisper-large-v3-"
     "turbo) writes new rows — merge, never erase", 14, bullet=True)
para(tb.text_frame, "align reconciles VTT timing + speaker attribution + STT by text alignment; "
     "millisecond offsets anchor every citation", 14, bullet=True)
para(tb.text_frame, "Moments segmented at a 20s gap — the p90 measured over 28 real drops / 7,983 turns",
     14, bullet=True)
notes(s, "1:15–1:25 — Encoding choices are measured, not guessed: chunk size, hybrid ratio, and the 20-second "
         "moment gap all come from measuring the real pulled corpus. Provided transcripts are never "
         "overwritten — the STT lane verifies alongside, which is what keeps citations stable.")

# ── 8 · Ingestion + derived objects ─────────────────────────────────────────
s = slide()
header(s, "Ingestion — video & transcripts in", "1:25–1:37")
picture(s, DIA / "ingest-seq.png", Inches(0.3), Inches(1.2), Inches(9.3), Inches(6.0))
tb = box(s, Inches(9.8), Inches(1.5), Inches(3.2), Inches(5.6))
para(tb.text_frame, "Derived objects", 18, bold=True, color=ACCENT, first=True, space_after=8)
para(tb.text_frame, "Aligned transcript segments", 14, bullet=True)
para(tb.text_frame, "Screens — durable identity across meetings", 14, bullet=True)
para(tb.text_frame, "Screenshots + frame OCR", 14, bullet=True)
para(tb.text_frame, "Moments — the citable spans", 14, bullet=True)
para(tb.text_frame, "Artifacts: ADRs & action items (LLM-extracted, human-approved before publish)",
     14, bullet=True)
notes(s, "1:25–1:37 — Everything enters through one door: a write-once drop plus POST /ingests. The worker "
         "walks eight checkpointed, idempotent stages — a crash resumes, a rerun overwrites only its own "
         "outputs. Out come the derived objects on the right; the LLM-extracted artifacts are the only "
         "model-written data, and a human approves them before they publish.")

# ── 9 · Data model ───────────────────────────────────────────────────────────
s = slide()
header(s, "The data model", "1:37–1:45")
picture(s, DIA / "erd.png", Inches(0.3), Inches(1.15), Inches(8.9), Inches(6.1))
tb = box(s, Inches(9.4), Inches(1.5), Inches(3.6), Inches(5.6))
para(tb.text_frame, "Why it holds up", 18, bold=True, color=ACCENT, first=True, space_after=8)
para(tb.text_frame, "MOMENT is the citation currency — minted once, never renumbered", 14, bullet=True)
para(tb.text_frame, "SCREEN and PARTICIPANT are durable cross-meeting identities — lineage of a slide "
     "or a person across the corpus", 14, bullet=True)
para(tb.text_frame, "JOB ↔ MEETING is 1:1 — re-processing re-arms the same job, so every citation "
     "survives re-ingestion", 14, bullet=True)
para(tb.text_frame, "Relational FK graph — traversable by SQL alone; Neo4j is a projection of it",
     14, bullet=True)
notes(s, "1:37–1:45 — The moment is the atom: a span of meeting timeline tying transcript segments to a "
         "screenshot. Screens and participants persist across meetings, which is what lets us ask for the "
         "lineage of a slide. And because a meeting keeps its job and its IDs forever, citations survive "
         "re-ingestion and augmentation.")

# ── 10 · A user asks a question ──────────────────────────────────────────────
s = slide()
header(s, "When a user asks a question", "1:45–2:00")
picture(s, DIA / "query-seq.png", Inches(0.3), Inches(1.2), Inches(9.7), Inches(5.5))
tb = box(s, Inches(10.2), Inches(1.5), Inches(2.9), Inches(5.4))
para(tb.text_frame, "Deterministic in, deterministic out", 16, bold=True, color=ACCENT, first=True,
     space_after=8)
para(tb.text_frame, "LLM touches only steps 7–8", 14, bullet=True)
para(tb.text_frame, "Validator rejects uncited claims — in code", 14, bullet=True)
para(tb.text_frame, "Live today: hybrid search + replay on this same path; cited chat is the Epic-3 build",
     14, bullet=True)
notes(s, "1:45–2:00 — The full Q&A path: classify to a traversal template, deterministic retrieval from "
         "both stores, LLM synthesis, then the citation validator resolves every marker against Postgres — "
         "or the answer dies. Search and replay run on this path today; the chat synthesis layer is the "
         "current epic.")

# ── 11 · Demo ────────────────────────────────────────────────────────────────
s = slide()
header(s, "Live demo", "2:00–3:00")
tb = box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.6))
steps = [
    ("1 · Ingested corpus (10s)", "Real Teams meetings, each ingested through the pipeline — "
     "stage-by-stage progress streamed live over SSE."),
    ("2 · Search (20s)", "One query across every meeting — hybrid keyword + semantic hits, "
     "highlighted snippets, each hit a moment."),
    ("3 · Replay (25s)", "Click a hit — the recording opens at that exact moment. "
     "This is what a citation resolves to."),
    ("4 · The point (5s)", "Not a chatbot over notes — an evidence system. "
     "Every answer is a moment you can watch."),
]
first = True
for head, body in steps:
    para(tb.text_frame, head, 20, bold=True, color=ACCENT, first=first, space_after=2)
    first = False
    para(tb.text_frame, body, 15, color=DARK, space_after=12)
tb = box(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7))
para(tb.text_frame, "http://127.0.0.1:5173", 16, color=GRAY, first=True, align=PP_ALIGN.CENTER)
notes(s, "2:00–3:00 — Switch to the browser. Full word-for-word script with fallbacks: demo-script.md. "
         "Beats: meetings list w/ live stage progress → search the rehearsed query → open the hit with a "
         "screenshot → video plays at startMs → close on the evidence line.")

# ── 12 · Close ───────────────────────────────────────────────────────────────
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.9), SLIDE_W, Inches(1.8))
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()
tb = box(s, Inches(1), Inches(3.3), Inches(11.3), Inches(1.2))
para(tb.text_frame, "Every answer is evidence you can replay.", 34, bold=True,
     color=RGBColor(0xFF, 0xFF, 0xFF), first=True, align=PP_ALIGN.CENTER)
tb = box(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6))
para(tb.text_frame, "MeetingMiner · Cameron Blake · AI Cohort 2026", 15, color=GRAY, first=True,
     align=PP_ALIGN.CENTER)
notes(s, "Hold during Q&A.")

prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)} slides")
